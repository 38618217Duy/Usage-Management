#!/usr/bin/env python3
"""
PPTX to Markdown Converter
Converts PowerPoint presentations to LLM-ready Markdown format.

Usage:
    python pptx-to-markdown.py <input.pptx> [output_dir]

Requirements:
    pip install python-pptx Pillow
"""

import sys
import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE


def sanitize_filename(name: str) -> str:
    """Sanitize string for use as filename."""
    return re.sub(r'[^\w\s-]', '', name).strip().replace(' ', '_')[:50]


def extract_text_from_shape(shape) -> str:
    """Extract text from a shape, handling different shape types."""
    texts = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            para_text = ""
            for run in paragraph.runs:
                para_text += run.text
            if para_text.strip():
                # Detect bullet level
                level = paragraph.level
                prefix = "  " * level + "- " if level > 0 else ""
                texts.append(f"{prefix}{para_text.strip()}")

    if shape.has_table:
        table = shape.table
        rows = []
        for row_idx, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
            if row_idx == 0:
                rows.append("|" + "|".join(["---"] * len(row.cells)) + "|")
        texts.append("\n".join(rows))

    return "\n".join(texts)


def extract_images(slide, slide_num: int, output_dir: Path) -> list:
    """Extract images from slide and save to output directory."""
    images_dir = output_dir / "images"
    images_dir.mkdir(exist_ok=True)

    extracted = []
    img_count = 0

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_count += 1
            image = shape.image

            # Determine extension
            ext = image.ext
            if ext == 'jpeg':
                ext = 'jpg'

            filename = f"slide_{slide_num:02d}_img_{img_count:02d}.{ext}"
            filepath = images_dir / filename

            with open(filepath, 'wb') as f:
                f.write(image.blob)

            extracted.append({
                'filename': filename,
                'path': f"images/{filename}",
                'width': shape.width,
                'height': shape.height
            })

    return extracted


def get_slide_title(slide) -> str:
    """Extract slide title."""
    if slide.shapes.title:
        return slide.shapes.title.text.strip()

    # Fallback: find first text shape with large font
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and len(text) < 100:
                return text

    return ""


def get_slide_notes(slide) -> str:
    """Extract speaker notes from slide."""
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame:
            return notes_frame.text.strip()
    return ""


def convert_pptx_to_markdown(pptx_path: str, output_dir: str = None) -> str:
    """
    Convert PPTX file to Markdown.

    Args:
        pptx_path: Path to the PPTX file
        output_dir: Output directory (default: same as input file)

    Returns:
        Path to the generated Markdown file
    """
    pptx_path = Path(pptx_path)

    if not pptx_path.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")

    if output_dir:
        output_path = Path(output_dir)
    else:
        output_path = pptx_path.parent / f"{pptx_path.stem}_output"

    output_path.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path)

    markdown_lines = []
    markdown_lines.append(f"# {pptx_path.stem}\n")
    markdown_lines.append(f"*Extracted from: {pptx_path.name}*\n")
    markdown_lines.append(f"*Total slides: {len(prs.slides)}*\n")
    markdown_lines.append("---\n")

    for slide_num, slide in enumerate(prs.slides, 1):
        # Slide header
        title = get_slide_title(slide)
        if title:
            markdown_lines.append(f"## Slide {slide_num}: {title}\n")
        else:
            markdown_lines.append(f"## Slide {slide_num}\n")

        # Extract content
        content_parts = []
        for shape in slide.shapes:
            # Skip title shape (already extracted)
            if shape == slide.shapes.title:
                continue

            text = extract_text_from_shape(shape)
            if text:
                content_parts.append(text)

        if content_parts:
            markdown_lines.append("\n".join(content_parts))
            markdown_lines.append("")

        # Extract and reference images
        images = extract_images(slide, slide_num, output_path)
        if images:
            markdown_lines.append("")
            for img in images:
                markdown_lines.append(f"![{img['filename']}]({img['path']})")
            markdown_lines.append("")

        # Speaker notes
        notes = get_slide_notes(slide)
        if notes:
            markdown_lines.append("")
            markdown_lines.append("> **Speaker Notes:**")
            for line in notes.split('\n'):
                if line.strip():
                    markdown_lines.append(f"> {line.strip()}")
            markdown_lines.append("")

        markdown_lines.append("\n---\n")

    # Write markdown file
    md_filename = f"{pptx_path.stem}.md"
    md_path = output_path / md_filename

    with open(md_path, 'w', encoding='utf-8') as f:
        f.write("\n".join(markdown_lines))

    print(f"Converted: {pptx_path.name}")
    print(f"Output: {md_path}")
    print(f"Images: {output_path / 'images'}")

    return str(md_path)


def batch_convert(input_dir: str, output_dir: str = None):
    """Convert all PPTX files in a directory."""
    input_path = Path(input_dir)
    pptx_files = list(input_path.glob("*.pptx"))

    if not pptx_files:
        print(f"No PPTX files found in {input_dir}")
        return

    print(f"Found {len(pptx_files)} PPTX files")

    for pptx_file in pptx_files:
        try:
            if output_dir:
                file_output = Path(output_dir) / pptx_file.stem
            else:
                file_output = None
            convert_pptx_to_markdown(str(pptx_file), file_output)
        except Exception as e:
            print(f"Error converting {pptx_file.name}: {e}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    input_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None

    if os.path.isdir(input_path):
        batch_convert(input_path, output_dir)
    else:
        convert_pptx_to_markdown(input_path, output_dir)
